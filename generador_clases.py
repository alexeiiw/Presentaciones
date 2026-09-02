from pathlib import Path
import random
import re

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from biblioteca_imagenes import buscar_en_biblioteca, guardar_en_biblioteca
from estilos import EstiloPresentacion, obtener_estilo
from parser_markdown import Clase, Diapositiva
from proveedores_imagenes import obtener_imagen, proveedores_desde_modo


ANCHO = Inches(13.333)
ALTO = Inches(7.5)


def generar_presentacion(
    clase: Clase,
    salida: Path,
    estilo_nombre: str | None = None,
    modo_imagen: str = "biblioteca_pexels_pixabay",
    distribucion: str = "alternada",
) -> Path:
    prs = Presentation()
    prs.slide_width = ANCHO
    prs.slide_height = ALTO
    estilo = obtener_estilo(estilo_nombre or clase.estilo)
    imagenes_usadas: set[str] = set()

    imagen_portada = _resolver_imagen(clase.imagen_universidad, modo_imagen, imagenes_usadas)
    logo_portada = _resolver_imagen(clase.logo_portada, modo_imagen, imagenes_usadas)
    _crear_portada(prs, clase, estilo, imagen_portada, logo_portada)
    if clase.agenda:
        _crear_diapositivas_lista(prs, "Agenda de la clase", clase.agenda, estilo, "Lo que veremos hoy", max_por_slide=7)

    contenido = clase.contenido_presentacion or [d.titulo for d in clase.diapositivas]
    if contenido:
        _crear_diapositivas_lista(prs, "Contenido de la presentacion", contenido, estilo, "Estructura del material", max_por_slide=10)

    for indice, diapositiva in enumerate(clase.diapositivas):
        imagen = _resolver_imagen(diapositiva.imagen, modo_imagen, imagenes_usadas)
        layout = _layout_para_diapositiva(indice, distribucion)

        if diapositiva.tipo == "codigo" or diapositiva.codigo:
            _crear_diapositiva_codigo(prs, diapositiva, estilo, imagen, modo_imagen, layout)
        elif diapositiva.tipo == "seccion" or _es_diapositiva_solo_titulo(diapositiva):
            _crear_diapositiva_seccion(prs, diapositiva, estilo, imagen, modo_imagen)
        elif diapositiva.tipo in {"columnas", "ruta", "frase", "diagrama", "actividad", "repositorio"}:
            _crear_diapositiva_avanzada(prs, diapositiva, estilo, imagen, modo_imagen)
        else:
            _crear_diapositiva_contenido(prs, diapositiva, estilo, imagen, modo_imagen, layout)

    aprendizajes = clase.aprendizajes or _aprendizajes_desde_diapositivas(clase.diapositivas)
    if aprendizajes or clase.frase_final:
        _crear_diapositiva_cierre(prs, aprendizajes, clase.frase_final, estilo)

    salida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(salida)
    return salida


def _crear_portada(prs: Presentation, clase: Clase, estilo: EstiloPresentacion, imagen: Path | None = None, logo: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)

    if imagen:
        _imagen_cover(slide, imagen, Inches(8.35), Inches(0), Inches(4.98), ALTO)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.35), Inches(0), Inches(4.98), ALTO)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = estilo.fondo_secundario
        overlay.fill.transparency = 0.28
        overlay.line.fill.background()
    else:
        _agregar_banda_visual(slide, estilo)

    if logo:
        _imagen_contain(slide, logo, Inches(0.82), Inches(0.55), Inches(1.6), Inches(0.95))

    tam_titulo = 38 if len(clase.titulo) > 55 else 42
    _texto(slide, clase.titulo, Inches(0.8), Inches(1.55), Inches(7.7), Inches(2.35), tam_titulo, estilo.titulo, estilo.fuente_titulo, negrita=True)
    subtitulo = " | ".join(x for x in [clase.asignatura, clase.universidad] if x)
    if subtitulo:
        _texto(slide, subtitulo, Inches(0.85), Inches(4.15), Inches(7.7), Inches(0.8), 17, estilo.texto, estilo.fuente_texto)
    if clase.profesor:
        _texto(slide, clase.profesor, Inches(0.85), Inches(6.05), Inches(6.4), Inches(0.4), 16, estilo.acento, estilo.fuente_texto)


def _crear_diapositivas_lista(prs: Presentation, titulo: str, items: list[str], estilo: EstiloPresentacion, subtitulo: str = "", max_por_slide: int = 10) -> None:
    for pagina, inicio in enumerate(range(0, len(items), max_por_slide), start=1):
        bloque = items[inicio : inicio + max_por_slide]
        titulo_pagina = titulo if len(items) <= max_por_slide else f"{titulo} ({pagina})"
        _crear_diapositiva_lista(prs, titulo_pagina, bloque, estilo, subtitulo, inicio)


def _crear_diapositiva_lista(prs: Presentation, titulo: str, items: list[str], estilo: EstiloPresentacion, subtitulo: str = "", offset: int = 0) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    _texto(slide, titulo, Inches(0.65), Inches(0.45), Inches(11.6), Inches(0.8), 32, estilo.titulo, estilo.fuente_titulo, negrita=True)
    _linea_acento(slide, estilo, Inches(0.65), Inches(1.35), Inches(2.0))
    if subtitulo:
        _texto(slide, subtitulo, Inches(0.7), Inches(1.55), Inches(10.8), Inches(0.4), 15, estilo.texto, estilo.fuente_texto)

    columnas = 2 if len(items) > 7 else 1
    ancho = Inches(5.5 if columnas == 2 else 10.8)
    alto = Inches(0.68 if columnas == 2 else 0.58)
    salto_y = Inches(0.86 if columnas == 2 else 0.72)
    for idx, item in enumerate(items):
        col = idx % columnas
        fila = idx // columnas
        x = Inches(0.8) + col * Inches(6.0)
        y = Inches(2.15) + fila * salto_y
        _tarjeta_item(slide, offset + idx + 1, item, estilo, x, y, ancho, alto)


def _crear_diapositiva_cierre(prs: Presentation, aprendizajes: list[str], frase: str, estilo: EstiloPresentacion) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    _agregar_banda_visual(slide, estilo)
    _texto(slide, "Esto fue lo que aprendiste", Inches(0.75), Inches(0.65), Inches(8.0), Inches(0.8), 34, estilo.titulo, estilo.fuente_titulo, negrita=True)
    _linea_acento(slide, estilo, Inches(0.75), Inches(1.55), Inches(2.0))
    _bullets(slide, aprendizajes[:5], Inches(1.0), Inches(2.0), Inches(7.5), Inches(3.0), estilo)

    frase = frase or "El conocimiento se consolida cuando puedes explicarlo, aplicarlo y mejorarlo."
    caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(8.0), Inches(0.95))
    caja.fill.solid()
    caja.fill.fore_color.rgb = estilo.caja
    caja.line.color.rgb = estilo.acento
    _texto(slide, frase, Inches(1.05), Inches(5.84), Inches(7.45), Inches(0.5), 17, estilo.acento, estilo.fuente_texto, negrita=True)


def _tarjeta_item(slide, numero: int, texto: str, estilo: EstiloPresentacion, x, y, w, h) -> None:
    caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    caja.fill.solid()
    caja.fill.fore_color.rgb = estilo.caja
    caja.line.color.rgb = estilo.acento
    tam = 12 if len(texto) > 70 else 14
    _texto(slide, f"{numero:02d}. {texto}", x + Inches(0.16), y + Inches(0.08), w - Inches(0.28), h - Inches(0.08), tam, estilo.texto, estilo.fuente_texto)


def _aprendizajes_desde_diapositivas(diapositivas: list[Diapositiva]) -> list[str]:
    return [f"Comprender {d.titulo.lower()}" for d in diapositivas[:5]]


def _resolver_imagen(keyword: str, modo_imagen: str, imagenes_usadas: set[str]) -> Path | None:
    if not keyword:
        return None

    local_directa = _resolver_ruta_local(keyword)
    if local_directa:
        imagenes_usadas.add(str(local_directa))
        return local_directa

    if modo_imagen in {"solo_diseno", "sin_imagen"}:
        return None

    if "biblioteca" in modo_imagen:
        local = buscar_en_biblioteca(keyword, imagenes_usadas)
        if local:
            imagenes_usadas.add(str(local))
            return local

    proveedores = proveedores_desde_modo(modo_imagen)
    if not proveedores:
        return None

    resultado = obtener_imagen(keyword, Path("temp") / "imagenes", proveedores, imagenes_usadas)
    if not resultado:
        return None

    ruta, proveedor, url = resultado
    guardada = guardar_en_biblioteca(keyword, ruta, proveedor, url)
    imagenes_usadas.add(str(guardada))
    return guardada


def _resolver_ruta_local(valor: str) -> Path | None:
    ruta = Path(valor.strip().strip('"'))
    extensiones = {".jpg", ".jpeg", ".png", ".webp", ".bmp"}
    if ruta.suffix.lower() not in extensiones:
        return None
    if ruta.exists() and ruta.is_file():
        return ruta
    ruta_assets = Path("assets") / ruta.name
    if ruta_assets.exists() and ruta_assets.is_file():
        return ruta_assets
    return None


def _layout_para_diapositiva(indice: int, distribucion: str) -> str:
    if distribucion == "fija":
        return "imagen_derecha"
    if distribucion == "aleatoria_controlada":
        return random.choice(["imagen_derecha", "imagen_izquierda"])
    return "imagen_izquierda" if indice % 2 else "imagen_derecha"


def _es_diapositiva_solo_titulo(d: Diapositiva) -> bool:
    return not d.objetivo and not d.contenido and not d.codigo and not d.imagen


def _crear_diapositiva_contenido(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion, imagen: Path | None, modo_imagen: str, layout: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)

    if layout == "imagen_izquierda":
        texto_x = Inches(5.85)
        imagen_x = Inches(0)
    else:
        texto_x = Inches(0.55)
        imagen_x = Inches(8.05)

    tam_titulo = 26 if len(d.titulo) > 42 else 30
    _texto(slide, d.titulo, texto_x, Inches(0.25), Inches(7.0), Inches(1.15), tam_titulo, estilo.titulo, estilo.fuente_titulo, negrita=True)
    _linea_acento(slide, estilo, texto_x, Inches(1.48), Inches(1.5))

    if imagen:
        _imagen_cover(slide, imagen, imagen_x, Inches(0.0), Inches(5.28), ALTO)
    elif modo_imagen in {"pexels_o_diseno", "solo_diseno", "biblioteca_pexels_pixabay", "biblioteca_pexels", "pexels_pixabay", "solo_biblioteca", "pixabay"}:
        _fallback_visual(slide, estilo, imagen_x, Inches(0.0), Inches(5.28), ALTO)

    y = Inches(1.72)
    if d.objetivo:
        _caja_objetivo(slide, d.objetivo, estilo, texto_x, y, Inches(6.9), Inches(0.9))
        y = Inches(2.82)

    _bullets(slide, d.contenido[:6], texto_x + Inches(0.2), y, Inches(6.55), Inches(3.95), estilo)


def _crear_diapositiva_codigo(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion, imagen: Path | None, modo_imagen: str, layout: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    tam_titulo = 25 if len(d.titulo) > 48 else 28
    _texto(slide, d.titulo, Inches(0.55), Inches(0.25), Inches(9.0), Inches(1.05), tam_titulo, estilo.titulo, estilo.fuente_titulo, negrita=True)
    _linea_acento(slide, estilo, Inches(0.55), Inches(1.42), Inches(1.5))

    if d.objetivo:
        _texto(slide, d.objetivo, Inches(0.6), Inches(1.58), Inches(5.8), Inches(0.55), 14, estilo.texto, estilo.fuente_texto)

    codigo = d.codigo or "\n".join(d.contenido)
    _bloque_codigo(slide, codigo, estilo, Inches(0.6), Inches(2.25), Inches(7.1), Inches(4.55))

    if imagen:
        _imagen_cover(slide, imagen, Inches(8.2), Inches(1.35), Inches(4.65), Inches(5.45))
    elif modo_imagen in {"pexels_o_diseno", "solo_diseno", "biblioteca_pexels_pixabay", "biblioteca_pexels", "pexels_pixabay", "solo_biblioteca", "pixabay"}:
        _fallback_visual(slide, estilo, Inches(8.2), Inches(1.35), Inches(4.65), Inches(5.45))


def _crear_diapositiva_avanzada(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion, imagen: Path | None, modo_imagen: str) -> None:
    if d.tipo == "columnas":
        _crear_diapositiva_columnas(prs, d, estilo)
    elif d.tipo == "ruta":
        _crear_diapositiva_ruta(prs, d, estilo)
    elif d.tipo == "frase":
        _crear_diapositiva_frase(prs, d, estilo, imagen, modo_imagen)
    elif d.tipo == "diagrama":
        _crear_diapositiva_diagrama(prs, d, estilo)
    elif d.tipo == "actividad":
        _crear_diapositiva_actividad(prs, d, estilo)
    elif d.tipo == "repositorio":
        _crear_diapositiva_repositorio(prs, d, estilo)


def _crear_diapositiva_seccion(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion, imagen: Path | None, modo_imagen: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    if imagen:
        _imagen_cover(slide, imagen, Inches(0), Inches(0), ANCHO, ALTO)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ANCHO, ALTO)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.42
        overlay.line.fill.background()
        color_titulo = RGBColor(255, 255, 255)
        color_texto = RGBColor(235, 235, 245)
    else:
        _agregar_banda_visual(slide, estilo)
        color_titulo = estilo.titulo
        color_texto = estilo.texto

    _texto(slide, "Nueva seccion", Inches(0.9), Inches(1.55), Inches(5.5), Inches(0.45), 16, estilo.acento, estilo.fuente_texto, negrita=True)
    _texto(slide, d.titulo, Inches(0.9), Inches(2.05), Inches(8.8), Inches(1.7), 38, color_titulo, estilo.fuente_titulo, negrita=True)
    subtitulo = d.objetivo or "Conectemos el bloque anterior con el siguiente tema tecnico."
    _texto(slide, subtitulo, Inches(0.95), Inches(4.15), Inches(7.8), Inches(0.75), 18, color_texto, estilo.fuente_texto)


def _crear_diapositiva_columnas(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    _encabezado(slide, d, estilo)
    items = d.contenido or [d.objetivo]
    mitad = max(1, (len(items) + 1) // 2)
    for col, bloque in enumerate([items[:mitad], items[mitad:]]):
        x = Inches(0.75 + col * 6.15)
        caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(5.55), Inches(4.95))
        caja.fill.solid()
        caja.fill.fore_color.rgb = estilo.caja
        caja.line.color.rgb = estilo.acento
        _texto(slide, "Idea clave" if col == 0 else "Aplicacion", x + Inches(0.25), Inches(2.05), Inches(5.0), Inches(0.45), 16, estilo.acento, estilo.fuente_texto, negrita=True)
        _bullets(slide, bloque[:5], x + Inches(0.35), Inches(2.65), Inches(4.85), Inches(3.55), estilo)


def _crear_diapositiva_ruta(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    _encabezado(slide, d, estilo)
    items = [item for item in (d.contenido or [d.objetivo]) if item][:5]
    for idx, item in enumerate(items):
        x = Inches(0.55 + idx * 2.55)
        y = Inches(2.55 if idx % 2 == 0 else 3.25)
        tarjeta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.78), Inches(2.15), Inches(2.15))
        tarjeta.fill.solid()
        tarjeta.fill.fore_color.rgb = estilo.caja
        tarjeta.line.color.rgb = estilo.acento
        circulo = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.82), Inches(0.82))
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = estilo.acento
        circulo.line.fill.background()
        _texto(slide, str(idx + 1), x + Inches(0.25), y + Inches(0.15), Inches(0.3), Inches(0.3), 16, RGBColor(255, 255, 255), estilo.fuente_texto, negrita=True)
        _texto_ajustado(slide, item, x + Inches(0.18), y + Inches(1.05), Inches(1.78), Inches(1.55), estilo.texto, estilo.fuente_texto, base=12)
        if idx < len(items) - 1:
            _linea_acento(slide, estilo, x + Inches(0.85), y + Inches(0.38), Inches(1.45))


def _crear_diapositiva_frase(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion, imagen: Path | None, modo_imagen: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    if imagen:
        _imagen_cover(slide, imagen, Inches(0), Inches(0), ANCHO, ALTO)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ANCHO, ALTO)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.35
        overlay.line.fill.background()
    elif modo_imagen != "sin_imagen":
        _agregar_banda_visual(slide, estilo)
    frase = d.objetivo or (d.contenido[0] if d.contenido else d.titulo)
    color_frase = RGBColor(255, 255, 255) if imagen else estilo.titulo
    _texto(slide, frase, Inches(1.2), Inches(2.15), Inches(10.6), Inches(2.0), 34, color_frase, estilo.fuente_titulo, negrita=True)
    if d.contenido[1:]:
        _texto(slide, d.contenido[1], Inches(1.25), Inches(4.65), Inches(9.8), Inches(0.6), 17, estilo.acento, estilo.fuente_texto)


def _crear_diapositiva_diagrama(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    _encabezado(slide, d, estilo)
    items = [item for item in (d.contenido or [d.objetivo]) if item][:6]
    centro_x, centro_y = Inches(4.75), Inches(3.15)
    centro = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, centro_x, centro_y, Inches(3.7), Inches(1.15))
    centro.fill.solid()
    centro.fill.fore_color.rgb = estilo.acento
    centro.line.fill.background()
    _texto_ajustado(slide, d.titulo, centro_x + Inches(0.22), centro_y + Inches(0.25), Inches(3.2), Inches(0.52), RGBColor(255, 255, 255), estilo.fuente_texto, base=15, negrita=True)
    posiciones = [(0.65, 1.78), (0.65, 3.36), (0.65, 4.94), (9.1, 1.78), (9.1, 3.36), (9.1, 4.94)]
    for idx, item in enumerate(items):
        x, y = posiciones[idx]
        caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.45), Inches(1.22))
        caja.fill.solid()
        caja.fill.fore_color.rgb = estilo.caja
        caja.line.color.rgb = estilo.acento
        _texto_ajustado(slide, item, Inches(x + 0.2), Inches(y + 0.16), Inches(3.0), Inches(0.72), estilo.texto, estilo.fuente_texto, base=11)


def _crear_diapositiva_actividad(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    _encabezado(slide, d, estilo)
    _texto(slide, d.objetivo or "Actividad guiada para aplicar el concepto en clase.", Inches(0.85), Inches(1.65), Inches(11.4), Inches(0.55), 17, estilo.acento, estilo.fuente_texto, negrita=True)
    etiquetas = ["Instrucciones", "Evidencia", "Cierre"]
    items = d.contenido or ["Resolver en equipos", "Compartir resultado", "Discutir hallazgos"]
    for idx, etiqueta in enumerate(etiquetas):
        x = Inches(0.85 + idx * 4.15)
        caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.55), Inches(3.55), Inches(3.4))
        caja.fill.solid()
        caja.fill.fore_color.rgb = estilo.caja
        caja.line.color.rgb = estilo.acento
        _texto(slide, etiqueta, x + Inches(0.25), Inches(2.85), Inches(3.0), Inches(0.45), 17, estilo.titulo, estilo.fuente_texto, negrita=True)
        _bullets(slide, items[idx::3][:3], x + Inches(0.3), Inches(3.45), Inches(2.95), Inches(1.8), estilo)


def _crear_diapositiva_repositorio(prs: Presentation, d: Diapositiva, estilo: EstiloPresentacion) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pintar_fondo(slide, estilo.fondo)
    _encabezado(slide, d, estilo)
    _texto(slide, d.objetivo or "Repositorio de recursos para continuar el aprendizaje.", Inches(0.85), Inches(1.65), Inches(11.4), Inches(0.5), 16, estilo.texto, estilo.fuente_texto)
    for idx, item in enumerate((d.contenido or [])[:6]):
        col = idx % 2
        fila = idx // 2
        _tarjeta_item(slide, idx + 1, item, estilo, Inches(0.85 + col * 6.0), Inches(2.45 + fila * 1.25), Inches(5.4), Inches(0.8))


def _pintar_fondo(slide, color: RGBColor) -> None:
    fondo = slide.background.fill
    fondo.solid()
    fondo.fore_color.rgb = color


def _encabezado(slide, d: Diapositiva, estilo: EstiloPresentacion) -> None:
    tam_titulo = 25 if len(d.titulo) > 52 else 30
    _texto(slide, d.titulo, Inches(0.65), Inches(0.35), Inches(11.8), Inches(0.9), tam_titulo, estilo.titulo, estilo.fuente_titulo, negrita=True)
    _linea_acento(slide, estilo, Inches(0.65), Inches(1.33), Inches(1.8))


def _agregar_banda_visual(slide, estilo: EstiloPresentacion) -> None:
    forma = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.1), 0, Inches(4.3), ALTO)
    forma.fill.solid()
    forma.fill.fore_color.rgb = estilo.fondo_secundario
    forma.line.fill.background()
    for i in range(5):
        barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.45 + i * 0.45), Inches(1.1 + i * 0.55), Inches(0.08), Inches(4.8 - i * 0.35))
        barra.fill.solid()
        barra.fill.fore_color.rgb = estilo.acento
        barra.line.fill.background()


def _linea_acento(slide, estilo: EstiloPresentacion, x, y, w) -> None:
    linea = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
    linea.fill.solid()
    linea.fill.fore_color.rgb = estilo.acento
    linea.line.fill.background()


def _texto(slide, texto: str, x, y, w, h, tam: int, color: RGBColor, fuente: str, negrita: bool = False) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _limpiar_texto(texto)
    p.font.size = Pt(tam)
    p.font.color.rgb = color
    p.font.name = fuente
    p.font.bold = negrita


def _texto_ajustado(slide, texto: str, x, y, w, h, color: RGBColor, fuente: str, base: int = 12, negrita: bool = False) -> None:
    tam = base
    if len(texto) > 130:
        tam = max(8, base - 3)
    elif len(texto) > 85:
        tam = max(9, base - 2)
    elif len(texto) > 55:
        tam = max(10, base - 1)
    _texto(slide, texto, x, y, w, h, tam, color, fuente, negrita=negrita)


def _caja_objetivo(slide, objetivo: str, estilo: EstiloPresentacion, x, y, w, h) -> None:
    caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    caja.fill.solid()
    caja.fill.fore_color.rgb = estilo.caja
    caja.line.color.rgb = estilo.acento
    _texto(slide, f"Objetivo: {objetivo}", x + Inches(0.18), y + Inches(0.12), w - Inches(0.35), h - Inches(0.2), 13, estilo.texto, estilo.fuente_texto)


def _bullets(slide, bullets: list[str], x, y, w, h, estilo: EstiloPresentacion) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = _limpiar_texto(item)
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = estilo.fuente_texto
        p.font.color.rgb = estilo.texto
        p.space_after = Pt(8)


def _bloque_codigo(slide, codigo: str, estilo: EstiloPresentacion, x, y, w, h) -> None:
    caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    caja.fill.solid()
    caja.fill.fore_color.rgb = estilo.codigo_fondo
    caja.line.color.rgb = estilo.acento
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), h - Inches(0.36))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _limpiar_texto(codigo[:1800])
    p.font.name = estilo.fuente_codigo
    p.font.size = Pt(13)
    p.font.color.rgb = estilo.codigo_texto


def _fallback_visual(slide, estilo: EstiloPresentacion, x, y, w, h) -> None:
    caja = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    caja.fill.solid()
    caja.fill.fore_color.rgb = estilo.fondo_secundario
    caja.line.fill.background()
    for i in range(6):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35 + i * 0.65), y + Inches(0.8 + i * 0.75), Inches(1.1), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = estilo.acento
        shape.fill.transparency = 0.45
        shape.line.fill.background()


def _imagen_cover(slide, ruta: Path, x, y, w, h) -> None:
    with Image.open(ruta) as img:
        img_w, img_h = img.size

    img_aspect = img_w / img_h
    box_aspect = w / h
    pic = slide.shapes.add_picture(str(ruta), x, y, width=w, height=h)

    if img_aspect > box_aspect:
        visible_w = box_aspect * img_h
        crop = max(0, (img_w - visible_w) / (2 * img_w))
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        visible_h = img_w / box_aspect
        crop = max(0, (img_h - visible_h) / (2 * img_h))
        pic.crop_top = crop
        pic.crop_bottom = crop


def _imagen_contain(slide, ruta: Path, x, y, w, h) -> None:
    with Image.open(ruta) as img:
        img_w, img_h = img.size

    escala = min(w / img_w, h / img_h)
    final_w = int(img_w * escala)
    final_h = int(img_h * escala)
    final_x = x + int((w - final_w) / 2)
    final_y = y + int((h - final_h) / 2)
    slide.shapes.add_picture(str(ruta), final_x, final_y, width=final_w, height=final_h)


def _limpiar_texto(texto: str) -> str:
    texto = texto.replace("\\(", "").replace("\\)", "")
    texto = texto.replace("**", "").replace("*", "")
    texto = re.sub(r"\\text\{([^}]*)\}", r"\1", texto)
    texto = re.sub(r"\\frac\{([^}]*)\}\{([^}]*)\}", r"\1 / \2", texto)
    texto = texto.replace("\\left", "").replace("\\right", "")
    texto = texto.replace("\\times", "x").replace("\\approx", "~")
    texto = texto.replace("\\eta", "eta").replace("\\mu", "micro")
    texto = texto.replace("\\", "")
    return texto
